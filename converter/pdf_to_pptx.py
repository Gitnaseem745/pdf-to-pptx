import sys
from pathlib import Path
import io
import os
import json
import time

# Get the project root directory (parent of converter/)
PROJECT_ROOT = Path(__file__).parent.parent
sys.path.append(str(PROJECT_ROOT))

# Load .env file if it exists
env_file = PROJECT_ROOT / ".env"
if env_file.exists():
    with open(env_file) as f:
        for line in f:
            line = line.strip()
            if line and not line.startswith('#') and '=' in line:
                key, value = line.split('=', 1)
                os.environ.setdefault(key.strip(), value.strip())

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    from google import genai
    from google.genai import types
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False
    print("⚠️  google-genai not installed. Run: pip3 install google-genai")

# Use absolute paths relative to the script location
INPUT_PDF = PROJECT_ROOT / "input" / "slides.pdf"
OUTPUT_PPTX = PROJECT_ROOT / "output" / "output.pptx"

# Gemini API configuration - Set your API key here or via environment variable
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")


def get_gemini_client():
    """Get Gemini API client with the API key."""
    if not GEMINI_API_KEY:
        print("❌ Error: GEMINI_API_KEY not set!")
        print("   Set it via environment variable: export GEMINI_API_KEY='your-api-key'")
        print("   Or edit the GEMINI_API_KEY variable in this script.")
        return None
    
    client = genai.Client(api_key=GEMINI_API_KEY)
    return client


def render_page_to_pil_image(page, dpi=100, for_ai=True):
    """Render PDF page to PIL Image.
    
    Args:
        page: PyMuPDF page object
        dpi: Resolution (100 for AI analysis, 200 for background image)
        for_ai: If True, resize for AI token efficiency. If False, keep high quality.
    """
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    
    # Only resize for AI analysis to save tokens
    if for_ai:
        max_dimension = 1200
        if img.width > max_dimension or img.height > max_dimension:
            ratio = min(max_dimension / img.width, max_dimension / img.height)
            new_size = (int(img.width * ratio), int(img.height * ratio))
            img = img.resize(new_size, Image.LANCZOS)
    
    return img


def analyze_slide_with_gemini(client, image, slide_num, max_retries=3):
    """Use Gemini AI to analyze a slide image and extract structured content."""
    
    # Convert PIL image to bytes
    img_bytes = io.BytesIO()
    image.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    image_data = img_bytes.read()
    
    prompt = """Analyze this presentation slide image and extract ALL content with precise styling information.

Return a JSON object with this EXACT structure:
{
    "background_color": "#HEXCOLOR",
    "elements": [
        {
            "type": "title" | "subtitle" | "heading" | "body" | "bullet" | "caption",
            "text": "exact text content",
            "position": {
                "x_percent": 0-100,
                "y_percent": 0-100,
                "width_percent": 0-100,
                "height_percent": 0-100
            },
            "style": {
                "font_size": number (in points, estimate based on visual size),
                "font_color": "#HEXCOLOR",
                "bold": true/false,
                "italic": true/false,
                "alignment": "left" | "center" | "right"
            },
            "bullet_level": 0-3 (0 = no bullet, 1+ = bullet indent level)
        }
    ]
}

IMPORTANT RULES:
1. Extract ALL text visible on the slide - don't miss anything
2. Preserve the EXACT text content including punctuation and special characters
3. Estimate positions as percentages of slide dimensions (0-100)
4. Identify text hierarchy: titles are usually larger at top, body text smaller below
5. Detect bullet points and their indent levels
6. Extract accurate colors - background and text colors (use hex format #RRGGBB)
7. Title/heading text is usually 28-48pt, body text 16-24pt, captions 10-14pt
8. For multi-line text blocks, include all lines in the "text" field separated by newlines
9. Return ONLY valid JSON, no markdown code blocks, no explanation"""

    # Models to try (fallback order) - prioritize by free tier RPM limits
    # gemini-2.5-flash-lite: 10 RPM (best for free tier)
    # gemini-2.5-flash: 5 RPM
    # gemini-3-flash: 5 RPM
    models = ["gemini-2.5-flash-lite", "gemini-2.5-flash", "gemini-2.0-flash"]
    
    for model_name in models:
        for attempt in range(max_retries):
            try:
                # Create image part for the API
                image_part = types.Part.from_bytes(data=image_data, mime_type="image/png")
                
                response = client.models.generate_content(
                    model=model_name,
                    contents=[prompt, image_part]
                )
                
                # Extract JSON from response
                response_text = response.text.strip()
                
                # Remove markdown code blocks if present
                if response_text.startswith("```json"):
                    response_text = response_text[7:]
                elif response_text.startswith("```"):
                    response_text = response_text[3:]
                if response_text.endswith("```"):
                    response_text = response_text[:-3]
                
                response_text = response_text.strip()
                
                slide_data = json.loads(response_text)
                element_count = len(slide_data.get('elements', []))
                print(f"   ✓ Slide {slide_num}: Extracted {element_count} elements")
                return slide_data
                
            except json.JSONDecodeError as e:
                print(f"   ⚠ Slide {slide_num}: JSON parse error - {e}")
                return None
            except Exception as e:
                error_str = str(e)
                if "429" in error_str or "RESOURCE_EXHAUSTED" in error_str:
                    # Rate limit - wait and retry with shorter intervals
                    # 10 RPM = 1 request per 6 seconds, so wait 10-20s should be enough
                    wait_time = 15 * (attempt + 1)  # 15s, 30s, 45s
                    print(f"   ⏳ Slide {slide_num}: Rate limited, waiting {wait_time}s (attempt {attempt + 1}/{max_retries})...")
                    time.sleep(wait_time)
                else:
                    print(f"   ⚠ Slide {slide_num}: AI error with {model_name} - {e}")
                    break  # Try next model
    
    print(f"   ⚠ Slide {slide_num}: All retries failed")
    return None


def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    if not hex_color:
        return (0, 0, 0)
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 6:
        try:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except ValueError:
            return (0, 0, 0)
    return (0, 0, 0)


def create_slide_from_ai_data(prs, slide_data, slide_width, slide_height, background_image=None):
    """Create a PPTX slide from AI-extracted data with optional background image."""
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Add background image first (if provided) - this preserves all visuals from PDF
    if background_image:
        img_bytes = io.BytesIO()
        background_image.save(img_bytes, format='PNG', optimize=True)
        img_bytes.seek(0)
        
        # Add image covering the entire slide as background
        bg_shape = slide.shapes.add_picture(
            img_bytes,
            Inches(0),
            Inches(0),
            width=slide_width,
            height=slide_height
        )
        # Send to back so text boxes appear on top
        # Move shape to beginning of shape tree (back)
        spTree = slide.shapes._spTree
        sp = bg_shape._element
        spTree.remove(sp)
        spTree.insert(2, sp)  # Insert after slide background elements
    else:
        # Set background color if no image
        bg_color = slide_data.get("background_color", "#FFFFFF")
        r, g, b = hex_to_rgb(bg_color)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
    
    # Process each element - add as editable text boxes on top of background
    elements = slide_data.get("elements", [])
    
    for elem in elements:
        text = elem.get("text", "").strip()
        if not text:
            continue
        
        pos = elem.get("position", {})
        style = elem.get("style", {})
        
        # Calculate position - use defaults if not provided
        x_pct = pos.get("x_percent", 5) / 100
        y_pct = pos.get("y_percent", 5) / 100
        w_pct = pos.get("width_percent", 90) / 100
        h_pct = pos.get("height_percent", 15) / 100
        
        # Convert percentages to slide coordinates
        left = int(x_pct * slide_width)
        top = int(y_pct * slide_height)
        width = int(w_pct * slide_width)
        height = int(h_pct * slide_height)
        
        # Ensure minimum dimensions
        if width < Inches(1):
            width = Inches(2)
        if height < Inches(0.4):
            height = Inches(0.6)
        
        # Clamp to slide boundaries
        if left < 0:
            left = int(Inches(0.2))
        if top < 0:
            top = int(Inches(0.2))
        if left + width > slide_width:
            width = int(slide_width - left - Inches(0.2))
        if top + height > slide_height:
            height = int(slide_height - top - Inches(0.2))
        
        # Create text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # Handle bullet points - split by newlines if multi-line
        lines = text.split('\n') if '\n' in text else [text]
        bullet_level = elem.get("bullet_level", 0)
        elem_type = elem.get("type", "body")
        
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            
            # Check if line starts with bullet character
            is_bullet_line = line.startswith('•') or line.startswith('-') or line.startswith('*')
            if is_bullet_line:
                line = line.lstrip('•-* ').strip()
                
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Set alignment
            alignment = style.get("alignment", "left")
            if alignment == "center":
                p.alignment = PP_ALIGN.CENTER
            elif alignment == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
            
            # Set bullet level for bullet items
            if bullet_level > 0 or is_bullet_line or elem_type == "bullet":
                p.level = max(0, (bullet_level or 1) - 1)
            
            # Add run with text
            run = p.add_run()
            run.text = line
            
            # Apply font styling
            font_size = style.get("font_size", 18)
            if isinstance(font_size, str):
                try:
                    font_size = int(font_size.replace('pt', ''))
                except:
                    font_size = 18
            
            # Adjust font size based on element type
            if elem_type == "title":
                font_size = max(font_size, 32)
            elif elem_type == "subtitle":
                font_size = max(font_size, 24)
            elif elem_type == "heading":
                font_size = max(font_size, 22)
            elif elem_type == "caption":
                font_size = min(font_size, 14)
            
            # Clamp font size to reasonable range
            font_size = max(10, min(60, font_size))
            run.font.size = Pt(font_size)
            
            # Font color
            font_color = style.get("font_color", "#000000")
            r, g, b = hex_to_rgb(font_color)
            run.font.color.rgb = RGBColor(r, g, b)
            
            # Bold and italic
            run.font.bold = style.get("bold", False) or elem_type in ["title", "heading"]
            run.font.italic = style.get("italic", False)
            
            # Font name
            run.font.name = "Arial"
    
    return slide


def create_pptx_from_pdf_with_ai(pdf_path, output_path):
    """Convert PDF to editable PPTX using Gemini AI for content extraction."""
    
    if not GEMINI_AVAILABLE:
        print("❌ google-genai package not available")
        return False
    
    client = get_gemini_client()
    if not client:
        return False
    
    print("🤖 Using Gemini AI for intelligent content extraction")
    print()
    
    # Open PDF
    doc = fitz.open(pdf_path)
    
    # Get first page dimensions to determine aspect ratio
    first_page = doc[0]
    pdf_width = first_page.rect.width
    pdf_height = first_page.rect.height
    aspect_ratio = pdf_width / pdf_height
    
    # Create presentation with matching aspect ratio
    prs = Presentation()
    
    if aspect_ratio > 1.5:  # Wide format (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:  # Standard format (4:3)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    total_pages = len(doc)
    print(f"📊 Processing {total_pages} slides...")
    print("   (Background image + editable text overlay)")
    print()
    
    failed_slides = []
    
    for page_num in range(total_pages):
        page = doc[page_num]
        
        # Render LOW-res image for AI analysis (saves tokens)
        ai_image = render_page_to_pil_image(page, dpi=100, for_ai=True)
        
        # Render HIGH-res image for background (preserves visual quality)
        bg_image = render_page_to_pil_image(page, dpi=200, for_ai=False)
        
        # Analyze with Gemini AI
        slide_data = analyze_slide_with_gemini(client, ai_image, page_num + 1)
        
        if slide_data and slide_data.get("elements"):
            # Create slide with background image + editable text overlay
            create_slide_from_ai_data(prs, slide_data, slide_width, slide_height, background_image=bg_image)
        else:
            # Fallback: add image as slide if AI fails
            failed_slides.append(page_num + 1)
            print(f"   ⚠ Slide {page_num + 1}: Using image fallback")
            
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            img_bytes = io.BytesIO()
            bg_image.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            
            slide.shapes.add_picture(
                img_bytes,
                Inches(0),
                Inches(0),
                width=slide_width,
                height=slide_height
            )
        
        # Rate limiting for free tier: gemini-2.5-flash-lite = 10 RPM
        # Using 10 seconds to stay well within limits and avoid retry loops
        if page_num < total_pages - 1:
            time.sleep(10)
    
    doc.close()
    prs.save(output_path)
    
    print()
    print(f"✅ PPTX created: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Editable slides: {len(prs.slides) - len(failed_slides)}")
    
    if failed_slides:
        print(f"   ⚠ Fallback slides (image only): {failed_slides}")
    
    return True


# ============ FALLBACK: Image-only mode (no AI) ============

def create_pptx_from_pdf_images(pdf_path, output_path):
    """Fallback: Convert PDF to PPTX using images only (not editable)."""
    
    doc = fitz.open(pdf_path)
    
    first_page = doc[0]
    pdf_width = first_page.rect.width
    pdf_height = first_page.rect.height
    aspect_ratio = pdf_width / pdf_height
    
    prs = Presentation()
    
    if aspect_ratio > 1.5:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    blank_layout = prs.slide_layouts[6]
    
    print("📷 Using image-only mode (slides not editable)")
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        slide = prs.slides.add_slide(blank_layout)
        
        # Render high-res image
        zoom = 200 / 72
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG', optimize=True)
        img_bytes.seek(0)
        
        slide.shapes.add_picture(
            img_bytes,
            Inches(0),
            Inches(0),
            width=slide_width,
            height=slide_height
        )
    
    doc.close()
    prs.save(output_path)
    
    print(f"✅ PPTX created: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print("   ⚠ Note: Slides contain images and are NOT editable")


def main():
    """Main entry point."""
    # Ensure output directory exists
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    
    print("=" * 60)
    print("📄 PDF to Editable PPTX Converter (with Gemini AI)")
    print("=" * 60)
    print()
    print(f"Input:  {INPUT_PDF}")
    print(f"Output: {OUTPUT_PPTX}")
    print()
    
    # Check if input file exists
    if not INPUT_PDF.exists():
        print(f"❌ Error: Input file not found: {INPUT_PDF}")
        return
    
    # Try AI-powered conversion first
    if GEMINI_AVAILABLE and GEMINI_API_KEY:
        success = create_pptx_from_pdf_with_ai(str(INPUT_PDF), str(OUTPUT_PPTX))
        if success:
            return
    
    # Fallback to image mode
    if not GEMINI_API_KEY:
        print("⚠ GEMINI_API_KEY not set - using image-only mode")
        print("   To enable editable slides, set: export GEMINI_API_KEY='your-key'")
        print()
    
    create_pptx_from_pdf_images(str(INPUT_PDF), str(OUTPUT_PPTX))


if __name__ == "__main__":
    main()
